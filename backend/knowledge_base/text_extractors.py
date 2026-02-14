"""Text extraction from various document formats.

References:
- Requirements 16: Document Processing
- Design Section 14.1: Processing Workflow
"""

import logging
import shutil
import subprocess
from abc import ABC, abstractmethod
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Optional

import markdown
import pandas as pd
import pdfplumber
import PyPDF2
from docx import Document as DocxDocument

try:
    from pptx import Presentation as PptxPresentation
except ImportError:  # pragma: no cover - tested via extractor runtime guard
    PptxPresentation = None

logger = logging.getLogger(__name__)


@dataclass
class ExtractionResult:
    """Result of text extraction from a document."""

    text: str
    metadata: Dict[str, Any]
    page_count: Optional[int] = None
    word_count: int = 0
    extraction_time: float = 0.0
    confidence: float = 1.0  # 1.0 for direct extraction, lower for OCR


class TextExtractor(ABC):
    """Abstract base class for text extractors."""

    @abstractmethod
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from a document.

        Args:
            file_path: Path to the document

        Returns:
            ExtractionResult with extracted text and metadata
        """
        pass

    def _count_words(self, text: str) -> int:
        """Count words in text.

        Args:
            text: Text to count words in

        Returns:
            Number of words
        """
        return len(text.split())


class PDFExtractor(TextExtractor):
    """Extract text from PDF files."""

    def __init__(self, use_pdfplumber: bool = True):
        """Initialize PDF extractor.

        Args:
            use_pdfplumber: Whether to use pdfplumber (better for tables) or PyPDF2
        """
        self.use_pdfplumber = use_pdfplumber

    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from PDF file.

        Args:
            file_path: Path to PDF file

        Returns:
            ExtractionResult with extracted text
        """
        start_time = datetime.now()

        try:
            if self.use_pdfplumber:
                text, metadata, page_count = self._extract_with_pdfplumber(file_path)
            else:
                text, metadata, page_count = self._extract_with_pypdf2(file_path)

            extraction_time = (datetime.now() - start_time).total_seconds()
            word_count = self._count_words(text)

            logger.info(
                "PDF extraction completed",
                extra={
                    "file": str(file_path),
                    "pages": page_count,
                    "words": word_count,
                    "time": extraction_time,
                },
            )

            return ExtractionResult(
                text=text,
                metadata=metadata,
                page_count=page_count,
                word_count=word_count,
                extraction_time=extraction_time,
            )

        except Exception as e:
            logger.error(f"PDF extraction failed: {e}", exc_info=True)
            raise

    def _extract_with_pdfplumber(self, file_path: Path) -> tuple[str, Dict, int]:
        """Extract text using pdfplumber."""
        text_parts = []
        metadata = {}

        with pdfplumber.open(file_path) as pdf:
            metadata = pdf.metadata or {}
            page_count = len(pdf.pages)

            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

        return "\n\n".join(text_parts), metadata, page_count

    def _extract_with_pypdf2(self, file_path: Path) -> tuple[str, Dict, int]:
        """Extract text using PyPDF2."""
        text_parts = []
        metadata = {}

        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            metadata = pdf_reader.metadata or {}
            page_count = len(pdf_reader.pages)

            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

        return "\n\n".join(text_parts), dict(metadata), page_count


class DOCXExtractor(TextExtractor):
    """Extract text from DOCX files."""

    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from DOCX file.

        Args:
            file_path: Path to DOCX file

        Returns:
            ExtractionResult with extracted text
        """
        start_time = datetime.now()

        try:
            doc = DocxDocument(file_path)

            # Extract text from paragraphs
            text_parts = [para.text for para in doc.paragraphs if para.text.strip()]

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)

            text = "\n\n".join(text_parts)

            # Extract metadata
            metadata = {
                "author": doc.core_properties.author,
                "title": doc.core_properties.title,
                "subject": doc.core_properties.subject,
                "created": (
                    str(doc.core_properties.created) if doc.core_properties.created else None
                ),
                "modified": (
                    str(doc.core_properties.modified) if doc.core_properties.modified else None
                ),
            }

            extraction_time = (datetime.now() - start_time).total_seconds()
            word_count = self._count_words(text)

            logger.info(
                "DOCX extraction completed",
                extra={
                    "file": str(file_path),
                    "words": word_count,
                    "time": extraction_time,
                },
            )

            return ExtractionResult(
                text=text,
                metadata=metadata,
                word_count=word_count,
                extraction_time=extraction_time,
            )

        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}", exc_info=True)
            raise


class DOCExtractor(TextExtractor):
    """Extract text from legacy DOC (binary Word) files."""

    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from legacy DOC file.

        Tries host-native converters first (textutil/antiword/catdoc), then
        attempts a DOCX parser fallback for mis-labelled files.
        """
        start_time = datetime.now()
        attempts: list[str] = []
        commands: list[tuple[str, list[str]]] = []

        if shutil.which("textutil"):
            commands.append(
                (
                    "textutil",
                    ["textutil", "-convert", "txt", "-stdout", str(file_path)],
                )
            )
        if shutil.which("antiword"):
            commands.append(("antiword", ["antiword", str(file_path)]))
        if shutil.which("catdoc"):
            commands.append(("catdoc", ["catdoc", str(file_path)]))

        for command_name, command in commands:
            try:
                completed = subprocess.run(
                    command,
                    check=True,
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    errors="ignore",
                )
                extracted_text = (completed.stdout or "").strip()
                if not extracted_text:
                    attempts.append(f"{command_name}: empty output")
                    continue

                extraction_time = (datetime.now() - start_time).total_seconds()
                word_count = self._count_words(extracted_text)
                return ExtractionResult(
                    text=extracted_text,
                    metadata={
                        "filename": file_path.name,
                        "size": file_path.stat().st_size,
                        "extractor": command_name,
                    },
                    word_count=word_count,
                    extraction_time=extraction_time,
                )
            except Exception as ex:
                attempts.append(f"{command_name}: {ex}")

        # Some clients mislabel DOCX as application/msword. Try DOCX parser last.
        try:
            fallback_result = DOCXExtractor().extract(file_path)
            fallback_result.metadata = {
                **fallback_result.metadata,
                "extractor": "python-docx-fallback",
            }
            return fallback_result
        except Exception as ex:
            attempts.append(f"python-docx-fallback: {ex}")

        attempts_summary = "; ".join(attempts) if attempts else "no extraction tool available"
        raise ValueError(
            "Legacy Word (.doc) extraction failed. "
            "Install textutil/antiword/catdoc or convert to .docx. "
            f"Details: {attempts_summary}"
        )


class PPTXExtractor(TextExtractor):
    """Extract text from PPTX slide decks."""

    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from PPTX file."""
        start_time = datetime.now()

        if PptxPresentation is None:
            raise ValueError(
                "PPT/PPTX extraction requires optional dependency `python-pptx`. "
                "Please install it and retry."
            )

        try:
            presentation = PptxPresentation(file_path)
            slide_count = len(presentation.slides)
            slide_text_parts: list[str] = []
            slide_summaries: list[dict[str, Any]] = []

            for slide_index, slide in enumerate(presentation.slides, start=1):
                lines: list[str] = []
                for shape in slide.shapes:
                    lines.extend(self._extract_shape_lines(shape))

                notes_text = self._extract_notes_text(slide)
                if notes_text:
                    lines.append(f"[Speaker Notes]\n{notes_text}")

                cleaned_lines = [line.strip() for line in lines if line and line.strip()]
                has_text = bool(cleaned_lines)
                if has_text:
                    slide_text = "\n".join(cleaned_lines)
                else:
                    slide_text = "(Empty slide)"

                slide_text_parts.append(f"[Slide {slide_index}]\n{slide_text}")
                slide_summaries.append(
                    {
                        "slide": slide_index,
                        "shape_count": len(slide.shapes),
                        "has_text": has_text,
                    }
                )

            text = "\n\n".join(slide_text_parts).strip()
            extraction_time = (datetime.now() - start_time).total_seconds()
            word_count = self._count_words(text)

            core = presentation.core_properties
            metadata = {
                "filename": file_path.name,
                "slide_count": slide_count,
                "slides": slide_summaries,
                "author": getattr(core, "author", None),
                "title": getattr(core, "title", None),
                "subject": getattr(core, "subject", None),
                "created": str(core.created) if getattr(core, "created", None) else None,
                "modified": str(core.modified) if getattr(core, "modified", None) else None,
            }

            logger.info(
                "PPTX extraction completed",
                extra={
                    "file": str(file_path),
                    "slides": slide_count,
                    "words": word_count,
                    "time": extraction_time,
                },
            )

            return ExtractionResult(
                text=text,
                metadata=metadata,
                page_count=slide_count,
                word_count=word_count,
                extraction_time=extraction_time,
            )
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}", exc_info=True)
            raise

    def _extract_shape_lines(self, shape: Any) -> list[str]:
        """Extract text from a single PowerPoint shape."""
        lines: list[str] = []

        if hasattr(shape, "shapes"):
            for child_shape in shape.shapes:
                lines.extend(self._extract_shape_lines(child_shape))

        if getattr(shape, "has_text_frame", False):
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is not None:
                frame_text = (getattr(text_frame, "text", "") or "").strip()
                if frame_text:
                    lines.extend(line.strip() for line in frame_text.splitlines() if line.strip())

        if getattr(shape, "has_table", False):
            table = getattr(shape, "table", None)
            if table is not None:
                for row in table.rows:
                    row_values = [cell.text.strip() for cell in row.cells]
                    if any(row_values):
                        lines.append(" | ".join(row_values))

        return lines

    @staticmethod
    def _extract_notes_text(slide: Any) -> str:
        """Extract speaker notes text from a slide."""
        if not getattr(slide, "has_notes_slide", False):
            return ""

        notes_slide = getattr(slide, "notes_slide", None)
        if notes_slide is None:
            return ""

        notes_frame = getattr(notes_slide, "notes_text_frame", None)
        if notes_frame is None:
            return ""

        return (getattr(notes_frame, "text", "") or "").strip()


class TextFileExtractor(TextExtractor):
    """Extract text from plain text files."""

    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from text file.

        Args:
            file_path: Path to text file

        Returns:
            ExtractionResult with extracted text
        """
        start_time = datetime.now()

        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()

            metadata = {
                "filename": file_path.name,
                "size": file_path.stat().st_size,
            }

            extraction_time = (datetime.now() - start_time).total_seconds()
            word_count = self._count_words(text)

            return ExtractionResult(
                text=text,
                metadata=metadata,
                word_count=word_count,
                extraction_time=extraction_time,
            )

        except Exception as e:
            logger.error(f"Text file extraction failed: {e}", exc_info=True)
            raise


class MarkdownExtractor(TextExtractor):
    """Extract text from Markdown files."""

    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from Markdown file.

        Args:
            file_path: Path to Markdown file

        Returns:
            ExtractionResult with extracted text (preserves formatting)
        """
        start_time = datetime.now()

        try:
            with open(file_path, "r", encoding="utf-8") as f:
                md_text = f.read()

            # Convert to HTML then extract text (preserves structure)
            html = markdown.markdown(md_text)

            # Keep both markdown and plain text
            metadata = {
                "filename": file_path.name,
                "size": file_path.stat().st_size,
                "format": "markdown",
                "html": html,
            }

            extraction_time = (datetime.now() - start_time).total_seconds()
            word_count = self._count_words(md_text)

            return ExtractionResult(
                text=md_text,  # Keep original markdown
                metadata=metadata,
                word_count=word_count,
                extraction_time=extraction_time,
            )

        except Exception as e:
            logger.error(f"Markdown extraction failed: {e}", exc_info=True)
            raise


class ExcelExtractor(TextExtractor):
    """Extract text from Excel spreadsheets (.xls/.xlsx)."""

    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from every worksheet as tab-separated lines."""
        start_time = datetime.now()

        try:
            # Read all sheets so we preserve workbook-level context.
            sheets = pd.read_excel(file_path, sheet_name=None, dtype=str)
            if not sheets:
                raise ValueError("Spreadsheet contains no worksheets")

            text_parts: list[str] = []
            sheet_summaries: list[dict[str, Any]] = []

            for sheet_name, dataframe in sheets.items():
                frame = dataframe.fillna("")
                if frame.columns.empty and frame.empty:
                    text_parts.append(f"[Sheet: {sheet_name}]\n(Empty sheet)")
                    sheet_summaries.append({"name": sheet_name, "rows": 0, "columns": 0})
                    continue

                header_line = "\t".join(str(col) for col in frame.columns)
                row_lines = [
                    "\t".join(str(cell) for cell in row)
                    for row in frame.itertuples(index=False, name=None)
                ]
                sheet_text = "\n".join([f"[Sheet: {sheet_name}]", header_line, *row_lines]).strip()
                text_parts.append(sheet_text)

                sheet_summaries.append(
                    {
                        "name": sheet_name,
                        "rows": int(frame.shape[0]),
                        "columns": int(frame.shape[1]),
                    }
                )

            text = "\n\n".join(part for part in text_parts if part.strip())
            extraction_time = (datetime.now() - start_time).total_seconds()
            word_count = self._count_words(text)

            metadata = {
                "filename": file_path.name,
                "sheet_count": len(sheet_summaries),
                "sheets": sheet_summaries,
            }

            logger.info(
                "Excel extraction completed",
                extra={
                    "file": str(file_path),
                    "sheets": len(sheet_summaries),
                    "words": word_count,
                    "time": extraction_time,
                },
            )

            return ExtractionResult(
                text=text,
                metadata=metadata,
                page_count=len(sheet_summaries),
                word_count=word_count,
                extraction_time=extraction_time,
            )

        except ImportError as e:
            logger.error(f"Excel extraction dependency missing: {e}", exc_info=True)
            raise ValueError(
                "Excel extraction requires optional dependencies (e.g. openpyxl/xlrd). "
                "Please install them and retry."
            ) from e
        except Exception as e:
            logger.error(f"Excel extraction failed: {e}", exc_info=True)
            raise


def get_extractor(file_type: str) -> TextExtractor:
    """Get appropriate text extractor for file type.

    Args:
        file_type: MIME type or file extension

    Returns:
        TextExtractor instance

    Raises:
        ValueError: If file type is not supported
    """
    normalized = file_type.lower()

    if "pdf" in normalized:
        return PDFExtractor()
    elif "application/msword" in normalized or normalized.endswith(".doc"):
        return DOCExtractor()
    elif "wordprocessingml.document" in normalized or "docx" in normalized or "word" in normalized:
        return DOCXExtractor()
    elif "presentationml.presentation" in normalized or normalized.endswith(".pptx"):
        return PPTXExtractor()
    elif (
        "spreadsheetml.sheet" in normalized
        or "vnd.ms-excel" in normalized
        or normalized.endswith(".xlsx")
        or normalized.endswith(".xls")
    ):
        return ExcelExtractor()
    elif "markdown" in normalized or normalized.endswith(".md"):
        return MarkdownExtractor()
    elif "text" in normalized or normalized.endswith(".txt"):
        return TextFileExtractor()
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
